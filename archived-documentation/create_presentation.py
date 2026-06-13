#!/usr/bin/env python3
"""
DSA Platform Presentation Generator
Creates PowerPoint presentation using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_image_if_exists(slide, image_path, left, top, width, height):
    """Add image to slide if file exists"""
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
            return True
        except Exception as e:
            print(f"  ⚠️  Could not add image: {os.path.basename(image_path)}")
            return False
    return False

def create_presentation():
    """Create DSA Platform presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    bg_color = RGBColor(10, 10, 10)  # Dark background
    primary_color = RGBColor(16, 185, 129)  # Emerald green
    text_color = RGBColor(255, 255, 255)  # White
    
    print("🚀 Creating DSA Platform Presentation...")
    print("=" * 50)
    
    # Slide 1: Title Slide
    print("Creating Slide 1: Title Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "🚀 Syntaxia"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "DSA Learning Platform\nFull-Stack EdTech Solution"
    for para in subtitle_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(28)
        para.font.color.rgb = text_color
    
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Presented By: Tushar\nDate: January 2025"
    for para in footer_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = text_color
    
    # Slide 2: Project Overview
    print("Creating Slide 2: Project Overview")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "📊 Project Overview", primary_color)
    
    content = """• 5 weeks development time
• DSA Problem Platform with scalability
• Next.js 16.1.6 + TypeScript 5.0 + PostgreSQL 16
• Full-Stack Web Application
• Admin Panel + User Interface"""
    
    # Add content on left side
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = content
    content_frame.word_wrap = True
    for para in content_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = text_color
        para.space_after = Pt(8)
    
    # Add HomePage screenshot on right side
    image_added = add_image_if_exists(slide, 'public/screenshots/HomePage.png', 
                                       Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.8))
    if not image_added:
        # Add placeholder text if image not found
        placeholder = slide.shapes.add_textbox(Inches(5.5), Inches(3.5), Inches(3.5), Inches(1))
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = "[Homepage Screenshot]"
        placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        placeholder_frame.paragraphs[0].font.size = Pt(16)
        placeholder_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    
    # Slide 3: Problem Statement
    print("Creating Slide 3: Problem Statement")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "❓ Problem Statement", primary_color)
    add_subtitle(slide, "Project Motivation:", primary_color)
    
    content = """Motivation:

• 🇧🇩 Build a DSA learning platform for Bangladesh
• 🎯 Enable students to learn and practice DSA
• 📚 Inspired by platforms like Algomaster
• 🔧 Create a customizable solution with admin control

Key Focus Areas:
• Admin panel for content management
• Visual progress tracking for students
• Clean, beginner-friendly interface"""
    
    add_content(slide, content, text_color, Inches(2))
    
    # Slide 4: Solution Approach
    print("Creating Slide 4: Solution Approach")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "💡 My Solution", primary_color)
    add_subtitle(slide, "What I Built:", primary_color)
    
    content = """✅ Learning Platform - DSA practice for students
✅ Admin Dashboard - Content management + Analytics
✅ Activity Heatmap - GitHub-style progress tracking
✅ Modern Clean UI - Beginner-friendly interface
✅ Clerk Authentication - Secure user + role management
✅ Monaco Editor - Professional code editor integration"""
    
    add_content(slide, content, text_color, Inches(2.2))
    
    # Slide 5: Tech Stack
    print("Creating Slide 5: Tech Stack")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🛠️ Tech Stack", primary_color)
    
    content = """Frontend:
• Next.js 16.1.6 (App Router) + TypeScript 5.0
• Tailwind CSS 3.4 + Monaco Editor + Radix UI

Backend:
• PostgreSQL 16 + Prisma ORM 7.4
• PrismaPg Adapter + Clerk Auth + Server Actions

Architecture:
• MVC Pattern + 3NF Database Design"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 6: System Architecture
    print("Creating Slide 6: System Architecture")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🏗️ System Architecture", primary_color)
    
    content = """Client (Next.js + React)
        ↓
Authentication (Clerk)
        ↓
Business Logic (Server Actions)
        ↓
Data Layer (Prisma ORM)
        ↓
Database (PostgreSQL - 12 Tables)

5 Sprints: Planning → Design → Implementation
→ Testing → Deployment"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 7: Key Features
    print("Creating Slide 7: Key Features")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "✨ Key Features", primary_color)
    
    content = """For Students:
• Browse problems with filters
• Monaco Editor + Real-time execution
• Activity Heatmap (365-day)
• Bookmarking & Progress tracking
• Celebration animations 🎉

For Admins:
• Complete CRUD Operations
• User Management + RBAC
• Analytics Dashboard
• Content Management"""
    
    # Add content on left side
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = content
    content_frame.word_wrap = True
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = text_color
        para.space_after = Pt(6)
    
    # Add hero demo screenshot on right side
    image_added = add_image_if_exists(slide, 'public/herodemo/herosec01.png', 
                                       Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.8))
    if not image_added:
        placeholder = slide.shapes.add_textbox(Inches(5.5), Inches(3.5), Inches(3.5), Inches(1))
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = "[Features Screenshot]"
        placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        placeholder_frame.paragraphs[0].font.size = Pt(16)
        placeholder_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    
    # Slide 8: Database Design
    print("Creating Slide 8: Database Design")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🗄️ Database Design", primary_color)
    
    content = """Complete 12 Tables in 3NF:
User • Problem • Solution • Topic • Tag • Company
TestCase • Bookmark • Discussion • Reply • Note • ActivityLog

Academic Justification:
• Third Normal Form (3NF) - Eliminates transitive
  dependencies and data redundancies

Key Relationships:
• User → Solution (1:M) • Problem → TestCase (1:M)
• Problem ↔ Tag (M:M) • Problem ↔ Company (M:M)

Performance Optimizations:
✅ B-tree indexes on foreign keys
✅ Connection pooling via PrismaPg adapter"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 9: Technical Challenges
    print("Creating Slide 9: Technical Challenges")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🔧 Technical Challenges", primary_color)
    
    content = """1. Code Execution Security ⚠️
   → Sandboxed environment, timeout limits
   ✅ Security measures implemented

2. Database Performance 🐌
   → Indexing + Connection pooling
   ✅ Query optimization applied

3. Activity Heatmap 📊
   → Client-side computation + Caching
   ✅ Rendering optimized"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 10: Testing Results
    print("Creating Slide 10: Testing Results")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🧪 Testing Results", primary_color)
    
    content = """Testing Coverage:
✅ 85%+ Code Coverage Achieved

Test Matrix:
✅ 50+ Unit Tests - All Passed
✅ 30+ Integration Tests - All Passed
✅ 10+ E2E System Workflows - All Passed
✅ Bug Resolution: 10/10 (100% Resolution Rate)

Security Audit:
✅ Security Score: 98/100
✅ SQL Injection - Prevented (Prisma parameterized)
✅ XSS - Prevented (React auto-escaping + CSP)
✅ Auth Bypass - Prevented (Clerk JWT + RBAC)"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 11: Performance Metrics
    print("Creating Slide 11: Performance Metrics")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "⚡ Performance & Results", primary_color)
    
    content = """Empirical Performance Data:
✅ Average API Response Time: 180ms
   (Target: <500ms - Exceeded by 64%)
✅ Lighthouse Performance Score: 96/100
✅ Database Query Time: <50ms average

Optimizations Implemented:
• Strategic B-tree indexing • Connection pooling
• Server Actions • Efficient queries • Data caching"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 12: Key Achievements
    print("Creating Slide 12: Key Achievements")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🏆 Key Achievements", primary_color)
    
    content = """Development Cycle:
✅ 5 weeks Agile development (5 sprints)
✅ 50+ React components built
✅ Maintainable, well-documented codebase

Technical Implementations:
✅ Complete CRUD operations for problems
✅ Clerk authentication with RBAC
✅ GitHub-style activity tracking (365-day)
✅ Real-time code execution engine
✅ Admin analytics dashboard

Performance Achievement:
✅ 96/100 Lighthouse Score
✅ 180ms average API response time"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 13: Live Demo
    print("Creating Slide 13: Live Demo")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🖥️ Live Demo Flow", primary_color)
    add_subtitle(slide, "Platform Walkthrough:", primary_color)
    
    content = """1. Browse Problems → Filter & Search
2. Code Editor → Write & Test
3. Submit Solution → Real-time feedback
4. Dashboard → Activity Heatmap
5. Admin Panel → CRUD Operations

Highlights: Monaco editor, Confetti 🎉, Heatmap"""
    
    # Add content on left side
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.5), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = content
    content_frame.word_wrap = True
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = text_color
        para.space_after = Pt(6)
    
    # Add CRUD Operation screenshot on right side
    image_added = add_image_if_exists(slide, 'public/screenshots/CRUD_Operation .png', 
                                       Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.8))
    if not image_added:
        # Try alternative filename
        image_added = add_image_if_exists(slide, 'public/screenshots/CRUD_Op.png', 
                                           Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.8))
    
    if not image_added:
        placeholder = slide.shapes.add_textbox(Inches(5.5), Inches(3.5), Inches(3.5), Inches(1))
        placeholder_frame = placeholder.text_frame
        placeholder_frame.text = "[Admin Panel Screenshot]"
        placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        placeholder_frame.paragraphs[0].font.size = Pt(16)
        placeholder_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
    
    # Slide 14: Future Roadmap
    print("Creating Slide 14: Future Roadmap")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🚀 Future Roadmap", primary_color)
    
    content = """Q2 2025:
• Multi-language support (Python, Java, C++)
• Enhanced search features
• AI-based difficulty calculator

Q3-Q4 2025:
• Contest system + Leaderboards
• Mobile app
• AI hints + Video tutorials
• Enterprise features"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Slide 15: Conclusion
    print("Creating Slide 15: Conclusion")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color
    
    add_title(slide, "🎯 Conclusion & Thank You", primary_color)
    
    content = """Project Summary:
✅ Addresses learning platform needs
✅ Full-stack web application
✅ Authentication with Clerk
✅ Database with Prisma ORM

Key Features:
DSA platform with admin panel, visual progress
tracking, and code editor integration

🙏 Thank You! Questions?

GitHub: @tusharsno"""
    
    add_content(slide, content, text_color, Inches(1.8))
    
    # Save presentation
    filename = 'DSA_Platform_Presentation.pptx'
    prs.save(filename)
    
    print("=" * 50)
    print(f"✅ Presentation created successfully!")
    print(f"📄 File: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("=" * 50)
    return filename

def add_title(slide, text, color):
    """Add title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = color

def add_subtitle(slide, text, color):
    """Add subtitle to slide"""
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = text
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.bold = True
    subtitle_para.font.color.rgb = color

def add_content(slide, text, color, top_margin):
    """Add content to slide"""
    content_box = slide.shapes.add_textbox(Inches(0.8), top_margin, Inches(8.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = text
    content_frame.word_wrap = True
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = color
        para.space_after = Pt(6)

if __name__ == "__main__":
    try:
        filename = create_presentation()
        print(f"\n🎉 Success! Open '{filename}' to view the presentation.")
    except ImportError:
        print("\n❌ Error: python-pptx library not found!")
        print("Install it using: pip install python-pptx")
    except Exception as e:
        print(f"\n❌ Error: {str(e)}")
